#!/usr/bin/env python3
"""
Medical PPT Generator for NanoClaw / DoctorNiuMa
Accepts JSON from stdin or first argument, outputs a .pptx file.

JSON schema:
{
  "title": str,
  "subtitle": str (optional),
  "author": str (optional),
  "date": str (optional),
  "style": "academic" | "clinical" | "patient" | "concise" | "bold_signal" |
           "dark_botanical" | "notebook" | "neon_cyber" | "swiss_modern" | "vintage",
  "output_path": str,  # e.g. "/workspace/group/report.pptx"
  "slides": [
    {
      "layout": "title" | "section" | "content" | "two_col" | "table" | "three_col" |
                "left_sidebar" | "right_sidebar" | "quote" | "center_focus" |
                "comparison" | "process" | "image_left" | "image_right" | "image_top" |
                "chart" | "big_number" | "split_panel" | "timeline" | "card_grid",
      "title": str,
      "content": [str, ...],     # bullet points (supports nested: "  - sub-bullet")
      "left": [str, ...],        # left column content
      "right": [str, ...],       # right column content
      "col1/col2/col3": [str],   # for three_col layout
      "sidebar": [str, ...],     # for sidebar layouts
      "table": [[str, ...], ...],# first row = header
      "attribution": str,        # for quote layout
      "left_label": str,         # for comparison layout
      "right_label": str,        # for comparison layout
      "notes": str,              # speaker notes
      # NotebookLM-inspired fields:
      "number": str,             # for big_number layout - the giant number/stat
      "label": str,              # for big_number layout - description under number
      "steps": [str, ...],       # for process/timeline layouts
      "cards": [{"title": str, "content": [str, ...]}, ...],  # for card_grid
      "image": {"path": str, "style": "plain" | "rounded" | "border"}  # for image layouts
    }
  ]
}
"""

import json
import sys
import os
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    import pptx.oxml.ns as nsmap
    from lxml import etree
    from PIL import Image as PILImage
    import os
    import sys
    from io import BytesIO
except ImportError as e:
    print(f"ERROR: Missing dependency: {e}", file=sys.stderr)
    sys.exit(1)


def create_chart_from_spec(chart_spec: dict, palette: dict):
    """Create chart image from specification (integrates chart_generator)."""
    try:
        # Import chart generator
        sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
        from chart_generator import create_chart

        return create_chart(chart_spec, palette)
    except ImportError:
        print("WARNING: Chart generator not available, skipping chart", file=sys.stderr)
        return None

# ─── Color Palette ────────────────────────────────────────────────────────────
PALETTES = {
    "academic": {
        "primary":        RGBColor(0x1B, 0x4F, 0x72),  # dark navy blue
        "secondary":      RGBColor(0x21, 0x8F, 0xBE),  # steel blue
        "accent":         RGBColor(0xE6, 0x7E, 0x22),  # amber
        "bg":             RGBColor(0xFF, 0xFF, 0xFF),
        "text":           RGBColor(0x1C, 0x2B, 0x3A),
        "subtext":        RGBColor(0x5D, 0x6D, 0x7E),
        "table_hdr":      RGBColor(0x1B, 0x4F, 0x72),
        "table_alt":      RGBColor(0xEB, 0xF5, 0xFB),
        # Gradient definitions
        "gradient_start": RGBColor(0x1B, 0x4F, 0x72),
        "gradient_end":   RGBColor(0x2C, 0x6F, 0x92),
        "background_tint": RGBColor(0xF0, 0xF5, 0xFA),
    },
    "clinical": {
        "primary":        RGBColor(0x0D, 0x47, 0xA1),  # clinical blue
        "secondary":      RGBColor(0x19, 0x76, 0xD2),
        "accent":         RGBColor(0xC6, 0x28, 0x28),  # alert red
        "bg":             RGBColor(0xFF, 0xFF, 0xFF),
        "text":           RGBColor(0x21, 0x21, 0x21),
        "subtext":        RGBColor(0x61, 0x61, 0x61),
        "table_hdr":      RGBColor(0x0D, 0x47, 0xA1),
        "table_alt":      RGBColor(0xE3, 0xF2, 0xFD),
        # Gradient definitions
        "gradient_start": RGBColor(0x0D, 0x47, 0xA1),
        "gradient_end":   RGBColor(0x1A, 0x60, 0xC0),
        "background_tint": RGBColor(0xE3, 0xF2, 0xFD),
    },
    "patient": {
        "primary":        RGBColor(0x00, 0x69, 0x5C),  # teal
        "secondary":      RGBColor(0x00, 0x96, 0x88),
        "accent":         RGBColor(0xFF, 0x8F, 0x00),  # warm amber
        "bg":             RGBColor(0xFF, 0xFF, 0xFF),
        "text":           RGBColor(0x21, 0x21, 0x21),
        "subtext":        RGBColor(0x54, 0x7F, 0x7A),
        "table_hdr":      RGBColor(0x00, 0x69, 0x5C),
        "table_alt":      RGBColor(0xE0, 0xF2, 0xF1),
        # Gradient definitions
        "gradient_start": RGBColor(0x00, 0x69, 0x5C),
        "gradient_end":   RGBColor(0x00, 0x80, 0x70),
        "background_tint": RGBColor(0xE0, 0xF2, 0xF1),
    },
    "concise": {
        "primary":        RGBColor(0x37, 0x47, 0x4F),  # slate
        "secondary":      RGBColor(0x54, 0x6E, 0x7A),
        "accent":         RGBColor(0x00, 0x89, 0x7B),  # teal accent
        "bg":             RGBColor(0xFF, 0xFF, 0xFF),
        "text":           RGBColor(0x21, 0x21, 0x21),
        "subtext":        RGBColor(0x60, 0x7D, 0x8B),
        "table_hdr":      RGBColor(0x37, 0x47, 0x4F),
        "table_alt":      RGBColor(0xEC, 0xEF, 0xF1),
        # Gradient definitions
        "gradient_start": RGBColor(0x37, 0x47, 0x4F),
        "gradient_end":   RGBColor(0x4A, 0x5A, 0x62),
        "background_tint": RGBColor(0xEC, 0xEF, 0xF1),
    },
    "bold_signal": {
        "primary":        RGBColor(0x00, 0xD4, 0xFF),  # electric blue
        "secondary":      RGBColor(0x2D, 0x2D, 0x2D),  # dark gray
        "accent":         RGBColor(0xFF, 0x2D, 0x78),  # hot pink
        "bg":             RGBColor(0x0D, 0x11, 0x17),  # dark bg
        "text":           RGBColor(0xFF, 0xFF, 0xFF),
        "subtext":        RGBColor(0xA0, 0xA8, 0xB4),
        "table_hdr":      RGBColor(0x00, 0xD4, 0xFF),
        "table_alt":      RGBColor(0x16, 0x1B, 0x22),
        "gradient_start": RGBColor(0x0D, 0x11, 0x17),
        "gradient_end":   RGBColor(0x1A, 0x1A, 0x2E),
        "background_tint": RGBColor(0x16, 0x1B, 0x22),
    },
    "dark_botanical": {
        "primary":        RGBColor(0x7C, 0x9A, 0x8E),  # sage
        "secondary":      RGBColor(0xD4, 0xA5, 0x74),  # warm accent
        "accent":         RGBColor(0xE8, 0xB4, 0xB8),  # pink accent
        "bg":             RGBColor(0x1A, 0x2F, 0x23),  # forest
        "text":           RGBColor(0xE8, 0xE4, 0xDF),  # cream text
        "subtext":        RGBColor(0x9A, 0x95, 0x90),
        "table_hdr":      RGBColor(0x7C, 0x9A, 0x8E),
        "table_alt":      RGBColor(0x22, 0x3B, 0x2D),
        "gradient_start": RGBColor(0x1A, 0x2F, 0x23),
        "gradient_end":   RGBColor(0x0F, 0x0F, 0x0F),
        "background_tint": RGBColor(0x22, 0x3B, 0x2D),
    },
    "notebook": {
        "primary":        RGBColor(0x6B, 0x6B, 0x6B),  # warm gray
        "secondary":      RGBColor(0xC6, 0x5D, 0x3E),  # terracotta accent
        "accent":         RGBColor(0xC6, 0x5D, 0x3E),  # terracotta
        "bg":             RGBColor(0xF5, 0xF1, 0xEB),  # off-white
        "text":           RGBColor(0x1A, 0x1A, 0x1A),
        "subtext":        RGBColor(0x6B, 0x6B, 0x6B),
        "table_hdr":      RGBColor(0x6B, 0x6B, 0x6B),
        "table_alt":      RGBColor(0xEE, 0xE8, 0xE0),
        "gradient_start": RGBColor(0xF5, 0xF1, 0xEB),
        "gradient_end":   RGBColor(0xE8, 0xE2, 0xD8),
        "background_tint": RGBColor(0xEE, 0xE8, 0xE0),
    },
    "neon_cyber": {
        "primary":        RGBColor(0x39, 0xFF, 0x14),  # neon green
        "secondary":      RGBColor(0x00, 0xFF, 0xFF),  # cyan
        "accent":         RGBColor(0xFF, 0x00, 0xAA),  # magenta
        "bg":             RGBColor(0x0A, 0x0A, 0x1A),  # deep purple
        "text":           RGBColor(0xE0, 0xE0, 0xE0),
        "subtext":        RGBColor(0x80, 0x80, 0x99),
        "table_hdr":      RGBColor(0x39, 0xFF, 0x14),
        "table_alt":      RGBColor(0x12, 0x12, 0x28),
        "gradient_start": RGBColor(0x0A, 0x0A, 0x1A),
        "gradient_end":   RGBColor(0x0A, 0x0F, 0x1C),
        "background_tint": RGBColor(0x12, 0x12, 0x28),
    },
    "swiss_modern": {
        "primary":        RGBColor(0xE3, 0x06, 0x13),  # red
        "secondary":      RGBColor(0x1A, 0x1A, 0x1A),  # black
        "accent":         RGBColor(0xE3, 0x06, 0x13),  # red
        "bg":             RGBColor(0xFF, 0xFF, 0xFF),  # white
        "text":           RGBColor(0x1A, 0x1A, 0x1A),
        "subtext":        RGBColor(0x55, 0x55, 0x55),
        "table_hdr":      RGBColor(0x1A, 0x1A, 0x1A),
        "table_alt":      RGBColor(0xF2, 0xF2, 0xF2),
        "gradient_start": RGBColor(0xFF, 0xFF, 0xFF),
        "gradient_end":   RGBColor(0xF5, 0xF5, 0xF5),
        "background_tint": RGBColor(0xF2, 0xF2, 0xF2),
    },
    "vintage": {
        "primary":        RGBColor(0x2D, 0x47, 0x39),  # deep green
        "secondary":      RGBColor(0xD4, 0xA8, 0x43),  # gold
        "accent":         RGBColor(0xD4, 0xA8, 0x43),  # gold
        "bg":             RGBColor(0xFA, 0xF3, 0xE8),  # cream
        "text":           RGBColor(0x2D, 0x47, 0x39),
        "subtext":        RGBColor(0x6B, 0x7F, 0x71),
        "table_hdr":      RGBColor(0x2D, 0x47, 0x39),
        "table_alt":      RGBColor(0xF0, 0xE8, 0xDA),
        "gradient_start": RGBColor(0xFA, 0xF3, 0xE8),
        "gradient_end":   RGBColor(0xF0, 0xE8, 0xDA),
        "background_tint": RGBColor(0xF0, 0xE8, 0xDA),
    },
    # NotebookLM-inspired styles
    "minimalist": {
        "primary":        RGBColor(0x11, 0x11, 0x11),  # jet black
        "secondary":      RGBColor(0x55, 0x55, 0x55),  # medium gray
        "accent":         RGBColor(0xFF, 0xCC, 0x00),  # electric yellow
        "bg":             RGBColor(0xF5, 0xF5, 0xF5),  # light gray
        "text":           RGBColor(0x11, 0x11, 0x11),
        "subtext":        RGBColor(0x66, 0x66, 0x66),
        "table_hdr":      RGBColor(0x11, 0x11, 0x11),
        "table_alt":      RGBColor(0xE9, 0xE9, 0xE9),
        "gradient_start": RGBColor(0xF5, 0xF5, 0xF5),
        "gradient_end":   RGBColor(0xE9, 0xE9, 0xE9),
        "background_tint": RGBColor(0xE9, 0xE9, 0xE9),
    },
    "editorial": {
        "primary":        RGBColor(0xFF, 0xCC, 0x00),  # high-sat yellow
        "secondary":      RGBColor(0x11, 0x11, 0x11),  # black
        "accent":         RGBColor(0xFF, 0x33, 0x33),  # alert red
        "bg":             RGBColor(0xFF, 0xFF, 0xFF),  # white
        "text":           RGBColor(0x11, 0x11, 0x11),
        "subtext":        RGBColor(0x44, 0x44, 0x44),
        "table_hdr":      RGBColor(0x11, 0x11, 0x11),
        "table_alt":      RGBColor(0xF2, 0xF2, 0xF2),
        "gradient_start": RGBColor(0xFF, 0xFF, 0xFF),
        "gradient_end":   RGBColor(0xF5, 0xF5, 0xF5),
        "background_tint": RGBColor(0xF2, 0xF2, 0xF2),
    },
    "neo_retro": {
        "primary":        RGBColor(0xFF, 0x00, 0x80),  # hot pink
        "secondary":      RGBColor(0x00, 0xCC, 0xCC),  # cyan
        "accent":         RGBColor(0xFF, 0xFF, 0x00),  # bright yellow
        "bg":             RGBColor(0xF5, 0xF0, 0xE1),  # cream grid-paper
        "text":           RGBColor(0x11, 0x11, 0x11),
        "subtext":        RGBColor(0x55, 0x55, 0x55),
        "table_hdr":      RGBColor(0x11, 0x11, 0x11),
        "table_alt":      RGBColor(0xE8, 0xE3, 0xD4),
        "gradient_start": RGBColor(0xF5, 0xF0, 0xE1),
        "gradient_end":   RGBColor(0xE8, 0xE3, 0xD4),
        "background_tint": RGBColor(0xE8, 0xE3, 0xD4),
    },
}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

def get_palette(style: str) -> dict:
    return PALETTES.get(style, PALETTES["academic"])


def set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_gradient_fill(slide, start_color: RGBColor, end_color: RGBColor, angle=5400000):
    """Add real gradient background to slide via XML manipulation.
    angle: rotation in 1/60000 degree units. 5400000 = 90° (top to bottom).
    0 = left to right, 2700000 = 45° diagonal."""
    bg = slide.background
    bgPr = bg._element
    # Remove existing fill
    for child in list(bgPr):
        if child.tag.endswith('}bgPr'):
            for fill_child in list(child):
                child.remove(fill_child)

    bgPr_elem = bgPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}bgPr')
    if bgPr_elem is None:
        # Create bgPr if it doesn't exist — fall back to solid
        set_bg(slide, start_color)
        return

    # Build gradFill XML
    nsmap_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    gradFill = etree.SubElement(bgPr_elem, f'{{{nsmap_a}}}gradFill')
    gradFill.set('rotWithShape', '1')
    gsLst = etree.SubElement(gradFill, f'{{{nsmap_a}}}gsLst')

    # Stop 1 (0%)
    gs1 = etree.SubElement(gsLst, f'{{{nsmap_a}}}gs')
    gs1.set('pos', '0')
    srgb1 = etree.SubElement(gs1, f'{{{nsmap_a}}}srgbClr')
    srgb1.set('val', f'{start_color.red:02X}{start_color.green:02X}{start_color.blue:02X}')

    # Stop 2 (100%)
    gs2 = etree.SubElement(gsLst, f'{{{nsmap_a}}}gs')
    gs2.set('pos', '100000')
    srgb2 = etree.SubElement(gs2, f'{{{nsmap_a}}}srgbClr')
    srgb2.set('val', f'{end_color.red:02X}{end_color.green:02X}{end_color.blue:02X}')

    # Linear direction
    lin = etree.SubElement(gradFill, f'{{{nsmap_a}}}lin')
    lin.set('ang', str(angle))
    lin.set('scaled', '1')


def add_translucent_shape(slide, shape_type, left, top, width, height,
                          fill_color: RGBColor, opacity_pct=30):
    """Add a semi-transparent decorative shape.
    shape_type: 1=rect, 9=oval, 5=rounded_rect"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()

    # Set transparency via XML (alpha = (100 - opacity_pct) * 1000)
    alpha_val = str((100 - opacity_pct) * 1000)
    nsmap_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    solidFill = shape.fill._fill
    srgbClr = solidFill.find(f'.//{{{nsmap_a}}}srgbClr')
    if srgbClr is not None:
        alpha_elem = etree.SubElement(srgbClr, f'{{{nsmap_a}}}alpha')
        alpha_elem.set('val', alpha_val)
    return shape


def add_side_bar(slide, color: RGBColor, width=Inches(0.3)):
    """Add colored side bar to slide."""
    add_rect(slide, Inches(0), Inches(0), width, SLIDE_H, color)


# Global style context — set in generate() and used by add_text_box/add_bullets_to_tf
_current_style = "academic"

def add_text_box(slide, text: str, left, top, width, height,
                 font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
                 wrap=True, font_name=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    # Apply font from style config
    fonts = FONTS.get(_current_style, FONTS["academic"])
    run.font.name = font_name or (fonts["heading"] if bold else fonts["body"])
    if color:
        run.font.color.rgb = color
    return txBox


def add_rect(slide, left, top, width, height, fill_color: RGBColor):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def parse_bullets(lines: list[str]) -> list[tuple[int, str]]:
    """Parse bullet lines into (level, text) tuples. Indent = 2 spaces per level."""
    result = []
    for line in lines:
        stripped = line.lstrip()
        indent = len(line) - len(stripped)
        level = indent // 2
        # Strip leading "- " or "• "
        text = stripped.lstrip('-•').strip()
        result.append((min(level, 2), text))
    return result


def add_bullets_to_tf(tf, bullets: list[tuple[int, str]], palette: dict,
                       base_size=18):
    sizes = [base_size, base_size - 2, base_size - 4]
    colors = [palette["text"], palette["subtext"], palette["subtext"]]
    fonts = FONTS.get(_current_style, FONTS["academic"])
    bullet_chars = ["●", "○", "–"]

    first = True
    for level, text in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = level
        p.space_before = Pt(4 if level == 0 else 2)
        p.space_after = Pt(2)
        # Add styled bullet character
        bullet = bullet_chars[min(level, 2)]
        run = p.add_run()
        run.text = f"{bullet}  {text}"
        run.font.size = Pt(sizes[level])
        run.font.color.rgb = colors[level]
        run.font.name = fonts["body"]


# ─── Typography ───────────────────────────────────────────────────────────────

FONTS = {
    "academic": {
        "heading": "Calibri",
        "body": "Calibri",
        "code": "Consolas",
        "size_h1": 36,
        "size_h2": 28,
        "size_h3": 24,
        "size_body": 20,
        "size_small": 16,
    },
    "clinical": {
        "heading": "Arial",
        "body": "Calibri",
        "code": "Consolas",
        "size_h1": 36,
        "size_h2": 28,
        "size_h3": 24,
        "size_body": 20,
        "size_small": 16,
    },
    "patient": {
        "heading": "Verdana",
        "body": "Calibri",
        "code": "Consolas",
        "size_h1": 36,
        "size_h2": 28,
        "size_h3": 24,
        "size_body": 20,
        "size_small": 16,
    },
    "concise": {
        "heading": "Calibri",
        "body": "Calibri",
        "code": "Consolas",
        "size_h1": 36,
        "size_h2": 28,
        "size_h3": 24,
        "size_body": 20,
        "size_small": 16,
    },
    "bold_signal": {
        "heading": "Archivo Black",
        "body": "Space Grotesk",
        "code": "Consolas",
        "size_h1": 40,
        "size_h2": 32,
        "size_h3": 26,
        "size_body": 20,
        "size_small": 16,
    },
    "dark_botanical": {
        "heading": "Cormorant",
        "body": "IBM Plex Sans",
        "code": "Consolas",
        "size_h1": 38,
        "size_h2": 30,
        "size_h3": 24,
        "size_body": 20,
        "size_small": 16,
    },
    "notebook": {
        "heading": "Bodoni MT",
        "body": "DM Sans",
        "code": "Consolas",
        "size_h1": 36,
        "size_h2": 28,
        "size_h3": 24,
        "size_body": 20,
        "size_small": 16,
    },
    "neon_cyber": {
        "heading": "Trebuchet MS",
        "body": "Calibri",
        "code": "Consolas",
        "size_h1": 40,
        "size_h2": 32,
        "size_h3": 26,
        "size_body": 20,
        "size_small": 16,
    },
    "swiss_modern": {
        "heading": "Arial Black",
        "body": "Arial",
        "code": "Consolas",
        "size_h1": 40,
        "size_h2": 32,
        "size_h3": 26,
        "size_body": 20,
        "size_small": 16,
    },
    "vintage": {
        "heading": "Georgia",
        "body": "Book Antiqua",
        "code": "Consolas",
        "size_h1": 36,
        "size_h2": 28,
        "size_h3": 24,
        "size_body": 20,
        "size_small": 16,
    },
    # NotebookLM-inspired fonts
    "minimalist": {
        "heading": "Impact",
        "body": "Helvetica",
        "code": "Consolas",
        "size_h1": 44,  # Extra large for NotebookLM dramatic headlines
        "size_h2": 32,
        "size_h3": 24,
        "size_body": 18,
        "size_small": 14,
    },
    "editorial": {
        "heading": "Impact",
        "body": "Helvetica",
        "code": "Consolas",
        "size_h1": 48,  # Maximum impact for editorial style
        "size_h2": 36,
        "size_h3": 26,
        "size_body": 18,
        "size_small": 14,
    },
    "neo_retro": {
        "heading": "Trebuchet MS",
        "body": "Verdana",
        "code": "Consolas",
        "size_h1": 40,
        "size_h2": 30,
        "size_h3": 24,
        "size_body": 18,
        "size_small": 14,
    },
}


def get_font_spec(style: str, element_type="body") -> tuple:
    """Get (font_name, font_size) for given style and element type."""
    fonts = FONTS.get(style, FONTS["academic"])
    element_type = element_type.lower()

    if element_type == "h1":
        return (fonts["heading"], fonts["size_h1"])
    elif element_type == "h2":
        return (fonts["heading"], fonts["size_h2"])
    elif element_type == "h3":
        return (fonts["heading"], fonts["size_h3"])
    elif element_type == "small":
        return (fonts["body"], fonts["size_small"])
    else:
        return (fonts["body"], fonts["size_body"])


# ─── Decorative Elements ─────────────────────────────────────────────────────

DECORATIVE_ELEMENTS = {
    "accent_bar": "Colored bar at slide edge",
    "corner_accent": "Triangle accent in corner",
    "circle_accent": "Partial circle border",
    "divider": "Horizontal line with style options",
    "icon_bullets": "Custom bullet icons (check, arrow, dot)",
    "number_badge": "Stylized numbers for lists",
}


def add_accent_bar(slide, position="left", color=None, palette=None):
    """Add accent bar at slide edge."""
    if palette and color is None:
        color = palette["primary"]
    if position == "left":
        add_rect(slide, Inches(0), Inches(0), Inches(0.3), SLIDE_H, color)
    elif position == "right":
        add_rect(slide, SLIDE_W - Inches(0.3), Inches(0), Inches(0.3), SLIDE_H, color)
    elif position == "top":
        add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), color)


def add_corner_accent(slide, corner="top-left", color=None, palette=None):
    """Add triangular accent in corner."""
    if palette and color is None:
        color = palette["accent"]
    size = Inches(1.5)
    if corner == "top-left":
        add_rect(slide, Inches(0), Inches(0), size, Inches(0.15), color)
        add_rect(slide, Inches(0), Inches(0), Inches(0.15), size, color)
    elif corner == "top-right":
        add_rect(slide, SLIDE_W - size, Inches(0), size, Inches(0.15), color)
        add_rect(slide, SLIDE_W - Inches(0.15), Inches(0), Inches(0.15), size, color)


def add_divider(slide, y_pos, width=Inches(10), color=None, palette=None, style="solid"):
    """Add horizontal divider line."""
    if palette and color is None:
        color = palette["subtext"]
    left = (SLIDE_W - width) / 2
    add_rect(slide, left, y_pos, width, Inches(0.05), color)


def add_decorations(slide, layout_type, palette):
    """Apply style-specific signature decorative elements."""
    style = _current_style

    if style == "bold_signal":
        _decorate_bold_signal(slide, layout_type, palette)
    elif style == "dark_botanical":
        _decorate_dark_botanical(slide, layout_type, palette)
    elif style == "notebook":
        _decorate_notebook(slide, layout_type, palette)
    elif style == "swiss_modern":
        _decorate_swiss_modern(slide, layout_type, palette)
    elif style == "neon_cyber":
        _decorate_neon_cyber(slide, layout_type, palette)
    elif style == "vintage":
        _decorate_vintage(slide, layout_type, palette)
    elif style == "minimalist":
        _decorate_minimalist(slide, layout_type, palette)
    elif style == "editorial":
        _decorate_editorial(slide, layout_type, palette)
    elif style == "neo_retro":
        _decorate_neo_retro(slide, layout_type, palette)
    else:
        # Default academic/clinical/patient/concise styles
        if layout_type in ["content", "two_col", "three_col", "left_sidebar",
                           "right_sidebar", "comparison"]:
            add_accent_bar(slide, position="left", palette=palette)
        if layout_type in ["quote", "center_focus"]:
            add_corner_accent(slide, corner="top-right", palette=palette)


def _decorate_minimalist(slide, layout_type, palette):
    """Minimalist: extreme whitespace, single accent line."""
    # Single thin accent line at bottom
    add_rect(slide, Inches(0.6), SLIDE_H - Inches(0.4), Inches(2), Inches(0.04), palette["accent"])
    # Very subtle gray corner mark (L-shape)
    add_rect(slide, Inches(0.4), Inches(0.4), Inches(0.03), Inches(0.6), palette["subtext"])
    add_rect(slide, Inches(0.4), Inches(0.4), Inches(0.6), Inches(0.03), palette["subtext"])


def _decorate_editorial(slide, layout_type, palette):
    """Editorial: magazine-style, yellow accent blocks."""
    # Yellow accent strip on left
    add_rect(slide, Inches(0), Inches(0), Inches(0.25), SLIDE_H, palette["primary"])
    # Small yellow corner accent
    add_rect(slide, SLIDE_W - Inches(1.5), Inches(0), Inches(1.5), Inches(0.12), palette["primary"])


def _decorate_neo_retro(slide, layout_type, palette):
    """Neo-Retro: grid-paper vibe with colorful blocks."""
    # Pink accent on left
    add_rect(slide, Inches(0), Inches(0), Inches(0.2), SLIDE_H, palette["primary"])
    # Cyan block at top right
    add_rect(slide, SLIDE_W - Inches(2), Inches(0), Inches(2), Inches(0.15), palette["secondary"])
    # Yellow accent at bottom
    add_rect(slide, Inches(0.2), SLIDE_H - Inches(0.12), Inches(3), Inches(0.12), palette["accent"])


def _decorate_bold_signal(slide, layout_type, palette):
    """Bold Signal: colored card focal point + large section numbers + glow accents."""
    # Bottom accent strip with primary color
    add_rect(slide, Inches(0), SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12), palette["primary"])
    # Translucent accent circle top-right
    add_translucent_shape(slide, 9, SLIDE_W - Inches(3), Inches(-0.5),
                          Inches(3.5), Inches(3.5), palette["primary"], opacity_pct=12)
    # Small accent dot bottom-left
    add_translucent_shape(slide, 9, Inches(0.3), SLIDE_H - Inches(1.5),
                          Inches(0.8), Inches(0.8), palette["accent"], opacity_pct=25)


def _decorate_dark_botanical(slide, layout_type, palette):
    """Dark Botanical: abstract gradient circles + thin vertical line separators."""
    # Large translucent circle background element
    add_translucent_shape(slide, 9, SLIDE_W - Inches(5), Inches(2),
                          Inches(6), Inches(6), palette["primary"], opacity_pct=10)
    # Smaller circle overlay
    add_translucent_shape(slide, 9, SLIDE_W - Inches(3.5), Inches(3.5),
                          Inches(3), Inches(3), palette["secondary"], opacity_pct=8)
    # Thin vertical accent line
    add_rect(slide, Inches(0.15), Inches(1.4), Inches(0.03), Inches(5.5), palette["secondary"])


def _decorate_notebook(slide, layout_type, palette):
    """Notebook: paper container + colorful section tabs."""
    if layout_type in ["content", "two_col", "three_col", "comparison", "left_sidebar", "right_sidebar"]:
        # Paper-like inner container with subtle shadow effect
        # Shadow (slightly offset darker rect)
        add_translucent_shape(slide, 1, Inches(0.45), Inches(1.45),
                              Inches(12.5), Inches(5.9), RGBColor(0x00, 0x00, 0x00), opacity_pct=8)
        # White paper container
        add_rect(slide, Inches(0.35), Inches(1.35), Inches(12.5), Inches(5.9),
                 RGBColor(0xFF, 0xFF, 0xFF))
    # Colorful tab indicators at top
    tab_colors = [palette["accent"], palette["primary"], palette["secondary"]]
    for i, tc in enumerate(tab_colors):
        add_rect(slide, Inches(1.0 + i * 1.2), Inches(1.18), Inches(0.8), Inches(0.22), tc)


def _decorate_swiss_modern(slide, layout_type, palette):
    """Swiss Modern: grid system + red accent lines."""
    # Bold red vertical accent line
    add_rect(slide, Inches(0.4), Inches(1.3), Inches(0.06), Inches(5.8), palette["accent"])
    # Horizontal red line at bottom
    add_rect(slide, Inches(0.4), SLIDE_H - Inches(0.6), Inches(3), Inches(0.04), palette["accent"])
    # Subtle grid dots (simulate with small squares)
    for row in range(3):
        for col in range(4):
            add_translucent_shape(slide, 1,
                                  Inches(10 + col * 0.5), Inches(2 + row * 1.5),
                                  Inches(0.04), Inches(0.04),
                                  palette["subtext"], opacity_pct=20)


def _decorate_neon_cyber(slide, layout_type, palette):
    """Neon Cyber: dark base + glowing border lines + accent shapes."""
    # Glowing top border line
    add_rect(slide, Inches(0), Inches(1.22), SLIDE_W, Inches(0.04), palette["primary"])
    # Side glow accent
    add_translucent_shape(slide, 1, SLIDE_W - Inches(0.08), Inches(1.3),
                          Inches(0.08), Inches(5.5), palette["secondary"], opacity_pct=40)
    # Corner glow blob
    add_translucent_shape(slide, 9, SLIDE_W - Inches(2.5), SLIDE_H - Inches(2.5),
                          Inches(3), Inches(3), palette["accent"], opacity_pct=8)


def _decorate_vintage(slide, layout_type, palette):
    """Vintage: double-line border frame + serif elegance."""
    # Outer border frame
    t = Inches(0.08)
    m = Inches(0.25)
    # Top lines
    add_rect(slide, m, m, SLIDE_W - 2*m, t, palette["secondary"])
    add_rect(slide, m + Inches(0.15), m + Inches(0.15),
             SLIDE_W - 2*m - Inches(0.3), Inches(0.03), palette["secondary"])
    # Bottom lines
    add_rect(slide, m, SLIDE_H - m - t, SLIDE_W - 2*m, t, palette["secondary"])
    add_rect(slide, m + Inches(0.15), SLIDE_H - m - t - Inches(0.15),
             SLIDE_W - 2*m - Inches(0.3), Inches(0.03), palette["secondary"])
    # Left lines
    add_rect(slide, m, m, t, SLIDE_H - 2*m, palette["secondary"])
    # Right lines
    add_rect(slide, SLIDE_W - m - t, m, t, SLIDE_H - 2*m, palette["secondary"])


def add_image_to_slide(slide, image_path: str, left, top, width, style="plain", palette=None):
    """Add an image to a slide with optional styling."""
    if not os.path.exists(image_path):
        print(f"WARNING: Image not found: {image_path}", file=sys.stderr)
        return None

    try:
        pic = slide.shapes.add_picture(image_path, left, top, width=width)
        pic = pic._element  # Get underlying element for advanced styling

        # Apply style effects
        if style == "rounded":
            # Rounded corners effect (simplified - creates border)
            pic.get_or_add_ln().get_or_add_spPr().get_or_add_ln().get_or_add_w()
        elif style == "shadow":
            # Shadow effect
            pass  # Would require XML manipulation
        elif style == "border" and palette:
            # Border effect
            pass  # Would require XML manipulation

        return pic
    except Exception as e:
        print(f"WARNING: Failed to add image {image_path}: {e}", file=sys.stderr)
        return None


def build_image_left_slide(prs, slide_data: dict, palette: dict):
    """Image on left (40%), text on right (60%)."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _apply_slide_bg(slide, palette)

    title = slide_data.get("title", "")
    _add_title_bar(slide, title, palette)

    # Image (left 40%)
    image_spec = slide_data.get("image", {})
    if isinstance(image_spec, str):
        image_spec = {"path": image_spec}
    image_path = image_spec.get("path", "")
    image_style = image_spec.get("style", "plain")

    if image_path and os.path.exists(image_path):
        img_width = Inches(4.8)
        add_image_to_slide(slide, image_path, Inches(0.5), Inches(1.4), img_width, image_style, palette)

    # Content (right 60%)
    content = slide_data.get("content", [])
    if content:
        txBox = slide.shapes.add_textbox(Inches(5.5), Inches(1.4), Inches(7.3), Inches(5.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        bullets = parse_bullets(content)
        add_bullets_to_tf(tf, bullets, palette, base_size=18)

    return slide


def build_image_right_slide(prs, slide_data: dict, palette: dict):
    """Text on left (60%), image on right (40%)."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _apply_slide_bg(slide, palette)

    title = slide_data.get("title", "")
    _add_title_bar(slide, title, palette)

    # Content (left 60%)
    content = slide_data.get("content", [])
    if content:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(7.3), Inches(5.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        bullets = parse_bullets(content)
        add_bullets_to_tf(tf, bullets, palette, base_size=18)

    # Image (right 40%)
    image_spec = slide_data.get("image", {})
    if isinstance(image_spec, str):
        image_spec = {"path": image_spec}
    image_path = image_spec.get("path", "")
    image_style = image_spec.get("style", "plain")

    if image_path and os.path.exists(image_path):
        img_width = Inches(4.8)
        add_image_to_slide(slide, image_path, Inches(8.0), Inches(1.5), img_width, image_style, palette)

    return slide


def build_image_top_slide(prs, slide_data: dict, palette: dict):
    """Image on top (50%), text on bottom (50%)."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _apply_slide_bg(slide, palette)

    title = slide_data.get("title", "")
    _add_title_bar(slide, title, palette)

    # Image (top 50% of content area)
    image_spec = slide_data.get("image", {})
    if isinstance(image_spec, str):
        image_spec = {"path": image_spec}
    image_path = image_spec.get("path", "")
    image_style = image_spec.get("style", "plain")

    if image_path and os.path.exists(image_path):
        img_width = Inches(10)
        add_image_to_slide(slide, image_path, Inches(1.5), Inches(1.5), img_width, image_style, palette)

    # Content (bottom 50% of content area)
    content = slide_data.get("content", [])
    if content:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        bullets = parse_bullets(content)
        add_bullets_to_tf(tf, bullets, palette, base_size=16)

    return slide


# ─── NotebookLM-Inspired Layouts ───────────────────────────────────────────────

def build_big_number_slide(prs, slide_data: dict, palette: dict):
    """NotebookLM-style: Giant number + small supporting text.
    Great for key statistics, metrics, or impactful data points."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _apply_slide_bg(slide, palette)

    style = _current_style
    fonts = FONTS.get(style, FONTS["academic"])
    is_dark = style in ("bold_signal", "dark_botanical", "neon_cyber")

    # Giant number (occupies 40% of slide width)
    big_number = slide_data.get("number", slide_data.get("title", ""))
    number_color = palette["primary"] if not is_dark else palette["text"]

    # Super large font for the number
    add_text_box(slide, big_number,
                 Inches(0.8), Inches(1.8), Inches(5), Inches(3.0),
                 font_size=96, bold=True, color=number_color, align=PP_ALIGN.LEFT)

    # Accent line under number
    add_rect(slide, Inches(0.8), Inches(4.6), Inches(2.5), Inches(0.06), palette["accent"])

    # Supporting text (smaller, on the right or below)
    label = slide_data.get("label", "")
    if label:
        add_text_box(slide, label,
                     Inches(0.8), Inches(4.9), Inches(11.5), Inches(1.0),
                     font_size=fonts["size_body"], color=palette["subtext"])

    content = slide_data.get("content", [])
    if content:
        txBox = slide.shapes.add_textbox(Inches(6.5), Inches(2.0), Inches(6.0), Inches(4.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        bullets = parse_bullets(content[:4])  # Max 4 supporting points
        add_bullets_to_tf(tf, bullets, palette, base_size=16)

    return slide


def build_split_panel_slide(prs, slide_data: dict, palette: dict):
    """NotebookLM-style: Asymmetric 70/30 or 60/40 split.
    Left: Visual/Image/Quote, Right: Dense content OR vice versa."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _apply_slide_bg(slide, palette)

    style = _current_style
    fonts = FONTS.get(style, FONTS["academic"])

    split_ratio = slide_data.get("split", "60-40")  # "60-40", "70-30", "50-50"
    if split_ratio == "70-30":
        left_width = 9.0
        right_start = 9.3
    elif split_ratio == "50-50":
        left_width = 6.3
        right_start = 6.6
    else:  # 60-40 default
        left_width = 7.8
        right_start = 8.1

    # Left panel (larger)
    left_content = slide_data.get("left", slide_data.get("content", []))
    title = slide_data.get("title", "")

    if title:
        add_text_box(slide, title,
                     Inches(0.6), Inches(0.3), Inches(left_width), Inches(1.0),
                     font_size=fonts["size_h2"], bold=True, color=palette["text"])

    if left_content:
        txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(left_width - 0.5), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        bullets = parse_bullets(left_content[:6])
        add_bullets_to_tf(tf, bullets, palette, base_size=18)

    # Vertical divider
    add_rect(slide, Inches(right_start - 0.15), Inches(0.8), Inches(0.03), Inches(6.0), palette["subtext"])

    # Right panel (smaller, sidebar style)
    right_content = slide_data.get("right", slide_data.get("sidebar", []))
    if right_content:
        # Optional accent bar at top of sidebar
        add_rect(slide, Inches(right_start), Inches(0.8), Inches(3.8), Inches(0.08), palette["accent"])

        txBox = slide.shapes.add_textbox(Inches(right_start), Inches(1.1), Inches(3.8), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        bullets = parse_bullets(right_content[:4])
        add_bullets_to_tf(tf, bullets, palette, base_size=14)

    return slide


def build_timeline_slide(prs, slide_data: dict, palette: dict):
    """NotebookLM-style: Vertical timeline with alternating left/right content."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _apply_slide_bg(slide, palette)

    style = _current_style
    fonts = FONTS.get(style, FONTS["academic"])

    title = slide_data.get("title", "")
    if title:
        add_text_box(slide, title,
                     Inches(0.6), Inches(0.2), Inches(12.1), Inches(0.8),
                     font_size=fonts["size_h2"], bold=True, color=palette["text"])

    # Central timeline spine
    spine_x = Inches(6.55)
    add_rect(slide, spine_x, Inches(1.2), Inches(0.05), Inches(5.8), palette["primary"])

    steps = slide_data.get("steps", slide_data.get("content", []))
    num_steps = min(len(steps), 6)

    for i, step in enumerate(steps[:6]):
        y_pos = Inches(1.4 + i * 0.9)

        # Circle node on timeline
        circle = slide.shapes.add_shape(9, spine_x - Inches(0.15), y_pos, Inches(0.35), Inches(0.35))
        circle.fill.solid()
        circle.fill.fore_color.rgb = palette["accent"]
        circle.line.fill.background()

        # Alternating left/right content
        is_left = i % 2 == 0
        if is_left:
            text_left = Inches(0.6)
            text_width = Inches(5.7)
        else:
            text_left = Inches(7.0)
            text_width = Inches(5.7)

        step_text = step if isinstance(step, str) else step[0] if step else ""
        add_text_box(slide, step_text,
                     text_left, y_pos, text_width, Inches(0.8),
                     font_size=14, color=palette["text"])

    return slide


def build_card_grid_slide(prs, slide_data: dict, palette: dict):
    """NotebookLM-style: 2-3 cards in a row, each with icon + title + description."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _apply_slide_bg(slide, palette)

    style = _current_style
    fonts = FONTS.get(style, FONTS["academic"])
    is_dark = style in ("bold_signal", "dark_botanical", "neon_cyber")

    title = slide_data.get("title", "")
    if title:
        add_text_box(slide, title,
                     Inches(0.6), Inches(0.2), Inches(12.1), Inches(0.8),
                     font_size=fonts["size_h2"], bold=True, color=palette["text"])

    cards = slide_data.get("cards", slide_data.get("content", []))
    if isinstance(cards[0], str) if cards else True:
        # Convert simple list to card objects
        cards = [{"title": c, "content": []} for c in cards]

    num_cards = min(len(cards), 3)
    card_width = Inches(4.0)
    card_height = Inches(5.0)
    gap = Inches(0.3)
    start_x = Inches(0.5)

    for i, card in enumerate(cards[:3]):
        x_pos = start_x + i * (card_width + gap)

        # Card background
        card_bg_color = palette["table_alt"] if not is_dark else palette["background_tint"]
        add_rect(slide, x_pos, Inches(1.3), card_width, card_height, card_bg_color)

        # Optional icon placeholder (circle)
        add_translucent_shape(slide, 9, x_pos + Inches(1.5), Inches(1.6),
                              Inches(1.0), Inches(1.0), palette["primary"], opacity_pct=30)

        # Card title
        card_title = card.get("title", card) if isinstance(card, dict) else card
        add_text_box(slide, card_title,
                     x_pos + Inches(0.2), Inches(2.8), card_width - Inches(0.4), Inches(0.8),
                     font_size=18, bold=True, color=palette["primary"],
                     align=PP_ALIGN.CENTER)

        # Card content
        card_content = card.get("content", []) if isinstance(card, dict) else []
        if card_content:
            txBox = slide.shapes.add_textbox(x_pos + Inches(0.2), Inches(3.6),
                                             card_width - Inches(0.4), Inches(2.4))
            tf = txBox.text_frame
            tf.word_wrap = True
            bullets = parse_bullets(card_content[:3])
            add_bullets_to_tf(tf, bullets, palette, base_size=12)

    return slide


# ─── Slide Builders ───────────────────────────────────────────────────────────

def build_title_slide(prs, data: dict, palette: dict):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    style = _current_style
    is_dark = style in ("bold_signal", "dark_botanical", "neon_cyber")
    is_notebooklm = style in ("minimalist", "editorial", "neo_retro")

    # Try gradient background first, fallback to solid
    try:
        add_gradient_fill(slide, palette["gradient_start"], palette["gradient_end"],
                          angle=5400000)
    except Exception:
        set_bg(slide, palette["bg"])

    fonts = FONTS.get(style, FONTS["academic"])

    if is_notebooklm:
        # NotebookLM-inspired: clean, editorial, dramatic
        if style == "editorial":
            # Yellow accent strip
            add_rect(slide, Inches(0), Inches(0), Inches(0.4), SLIDE_H, palette["primary"])
            # Giant title
            title_text = data.get("title", "")
            add_text_box(slide, title_text,
                         Inches(0.8), Inches(1.8), Inches(12), Inches(3.0),
                         font_size=fonts["size_h1"] + 8, bold=True,
                         color=palette["text"])
            # Thin divider
            add_rect(slide, Inches(0.8), Inches(4.8), Inches(4), Inches(0.04), palette["primary"])
        elif style == "minimalist":
            # Maximum whitespace
            title_text = data.get("title", "")
            add_text_box(slide, title_text,
                         Inches(0.8), Inches(2.2), Inches(11.5), Inches(2.5),
                         font_size=fonts["size_h1"] + 6, bold=True,
                         color=palette["text"])
            # Single yellow accent line
            add_rect(slide, Inches(0.8), Inches(4.8), Inches(2), Inches(0.06), palette["accent"])
        elif style == "neo_retro":
            # Colorful blocks
            add_rect(slide, Inches(0), Inches(0), Inches(0.3), SLIDE_H, palette["primary"])
            add_rect(slide, SLIDE_W - Inches(2.5), Inches(0), Inches(2.5), Inches(0.2), palette["secondary"])
            add_rect(slide, Inches(0.3), SLIDE_H - Inches(0.15), Inches(4), Inches(0.15), palette["accent"])
            title_text = data.get("title", "")
            add_text_box(slide, title_text,
                         Inches(0.8), Inches(2.0), Inches(11.5), Inches(2.5),
                         font_size=fonts["size_h1"] + 4, bold=True,
                         color=palette["text"])

        subtitle = data.get("subtitle", "")
        if subtitle:
            add_text_box(slide, subtitle,
                         Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.8),
                         font_size=fonts["size_body"], color=palette["subtext"])

    elif is_dark:
        # Dark themes: full-bleed gradient, centered dramatic title
        # Large translucent decorative circle
        add_translucent_shape(slide, 9, Inches(7), Inches(1), Inches(7), Inches(7),
                              palette["primary"], opacity_pct=15)

        title_text = data.get("title", "")
        add_text_box(slide, title_text,
                     Inches(1.0), Inches(2.0), Inches(11.3), Inches(2.2),
                     font_size=fonts["size_h1"] + 4, bold=True,
                     color=palette["text"], align=PP_ALIGN.LEFT)

        # Accent line under title
        add_rect(slide, Inches(1.0), Inches(4.0), Inches(3.0), Inches(0.06), palette["primary"])

        subtitle = data.get("subtitle", "")
        if subtitle:
            add_text_box(slide, subtitle,
                         Inches(1.0), Inches(4.3), Inches(11.3), Inches(1.0),
                         font_size=fonts["size_body"], color=palette["subtext"])
    elif style == "swiss_modern":
        # Swiss: asymmetric layout, big red accent
        set_bg(slide, palette["bg"])
        # Large red block
        add_rect(slide, Inches(0), Inches(0), Inches(5), SLIDE_H, palette["accent"])
        title_text = data.get("title", "")
        add_text_box(slide, title_text,
                     Inches(5.5), Inches(2.0), Inches(7.3), Inches(2.5),
                     font_size=fonts["size_h1"] + 6, bold=True,
                     color=palette["text"], align=PP_ALIGN.LEFT)
        subtitle = data.get("subtitle", "")
        if subtitle:
            add_text_box(slide, subtitle,
                         Inches(5.5), Inches(4.5), Inches(7.3), Inches(0.8),
                         font_size=fonts["size_body"], color=palette["subtext"])
    elif style == "notebook":
        # Notebook: warm paper bg with paper container
        set_bg(slide, palette["bg"])
        # White paper card
        add_translucent_shape(slide, 1, Inches(0.45), Inches(1.55),
                              Inches(12.5), Inches(5.5), RGBColor(0x00, 0x00, 0x00), opacity_pct=6)
        add_rect(slide, Inches(0.35), Inches(1.45), Inches(12.5), Inches(5.5),
                 RGBColor(0xFF, 0xFF, 0xFF))
        # Colorful tabs
        tab_colors = [palette["accent"], palette["primary"], RGBColor(0x4A, 0x90, 0xD9)]
        for i, tc in enumerate(tab_colors):
            add_rect(slide, Inches(1.0 + i * 1.2), Inches(1.28), Inches(0.8), Inches(0.22), tc)

        title_text = data.get("title", "")
        add_text_box(slide, title_text,
                     Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.0),
                     font_size=fonts["size_h1"], bold=True, color=palette["text"])
        subtitle = data.get("subtitle", "")
        if subtitle:
            add_text_box(slide, subtitle,
                         Inches(1.0), Inches(4.0), Inches(11.3), Inches(0.8),
                         font_size=fonts["size_body"], color=palette["subtext"])
    elif style == "vintage":
        set_bg(slide, palette["bg"])
        # Ornamental double border
        m = Inches(0.4)
        t = Inches(0.06)
        add_rect(slide, m, m, SLIDE_W - 2*m, t, palette["secondary"])
        add_rect(slide, m, SLIDE_H - m - t, SLIDE_W - 2*m, t, palette["secondary"])
        add_rect(slide, m, m, t, SLIDE_H - 2*m, palette["secondary"])
        add_rect(slide, SLIDE_W - m - t, m, t, SLIDE_H - 2*m, palette["secondary"])
        # Inner border
        m2 = Inches(0.6)
        add_rect(slide, m2, m2, SLIDE_W - 2*m2, Inches(0.03), palette["secondary"])
        add_rect(slide, m2, SLIDE_H - m2, SLIDE_W - 2*m2, Inches(0.03), palette["secondary"])
        add_rect(slide, m2, m2, Inches(0.03), SLIDE_H - 2*m2, palette["secondary"])
        add_rect(slide, SLIDE_W - m2, m2, Inches(0.03), SLIDE_H - 2*m2, palette["secondary"])

        title_text = data.get("title", "")
        add_text_box(slide, title_text,
                     Inches(1.2), Inches(2.2), Inches(10.9), Inches(2.0),
                     font_size=fonts["size_h1"] + 2, bold=True,
                     color=palette["text"], align=PP_ALIGN.CENTER)
        # Decorative divider
        add_rect(slide, Inches(4.5), Inches(4.2), Inches(4.3), Inches(0.04), palette["secondary"])
        subtitle = data.get("subtitle", "")
        if subtitle:
            add_text_box(slide, subtitle,
                         Inches(1.2), Inches(4.5), Inches(10.9), Inches(0.8),
                         font_size=fonts["size_body"], color=palette["subtext"],
                         align=PP_ALIGN.CENTER)
    else:
        # Default academic/clinical/patient/concise: clean, gradient title bar
        set_bg(slide, palette["bg"])
        # Gradient-like top block (primary color)
        add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(3.5), palette["primary"])
        # Accent line
        add_rect(slide, Inches(0), Inches(3.5), SLIDE_W, Inches(0.06), palette["secondary"])

        title_text = data.get("title", "")
        add_text_box(slide, title_text,
                     Inches(1.0), Inches(1.0), Inches(11.3), Inches(2.0),
                     font_size=fonts["size_h1"] + 2, bold=True,
                     color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.LEFT)

        subtitle = data.get("subtitle", "")
        if subtitle:
            add_text_box(slide, subtitle,
                         Inches(1.0), Inches(4.0), Inches(11.3), Inches(1.0),
                         font_size=fonts["size_h2"] - 4, color=palette["subtext"])

    # Author + date (always at bottom)
    meta_parts = []
    if data.get("author"):
        meta_parts.append(data["author"])
    if data.get("date"):
        meta_parts.append(data["date"])
    else:
        meta_parts.append(datetime.now().strftime("%Y-%m-%d"))
    meta = "   |   ".join(meta_parts)
    meta_color = palette["subtext"] if not is_dark else palette.get("subtext", RGBColor(0xA0, 0xA8, 0xB4))
    add_text_box(slide, meta,
                 Inches(1.0), Inches(6.5), Inches(11.3), Inches(0.5),
                 font_size=14, color=meta_color)
    return slide


def build_section_slide(prs, slide_data: dict, palette: dict):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    style = _current_style
    fonts = FONTS.get(style, FONTS["academic"])
    is_dark = style in ("bold_signal", "dark_botanical", "neon_cyber")

    if is_dark:
        try:
            add_gradient_fill(slide, palette["gradient_start"], palette["gradient_end"], angle=2700000)
        except Exception:
            set_bg(slide, palette["bg"])
        # Large decorative accent
        add_translucent_shape(slide, 9, Inches(-2), Inches(1), Inches(6), Inches(6),
                              palette["primary"], opacity_pct=12)
    elif style == "swiss_modern":
        set_bg(slide, palette["bg"])
        add_rect(slide, Inches(0), Inches(0), Inches(2), SLIDE_H, palette["accent"])
    elif style == "vintage":
        set_bg(slide, palette["bg"])
    else:
        set_bg(slide, palette["primary"])

    title = slide_data.get("title", "")
    title_color = palette["text"] if is_dark else RGBColor(0xFF, 0xFF, 0xFF)
    if style == "swiss_modern":
        title_color = palette["text"]
    elif style == "vintage":
        title_color = palette["text"]

    add_text_box(slide, title,
                 Inches(1.5), Inches(2.5), Inches(10), Inches(2.0),
                 font_size=fonts["size_h1"] + 4, bold=True,
                 color=title_color, align=PP_ALIGN.CENTER if style not in ("swiss_modern",) else PP_ALIGN.LEFT)

    # Accent line under title
    if style == "swiss_modern":
        add_rect(slide, Inches(2.5), Inches(4.5), Inches(3), Inches(0.06), palette["accent"])
    elif style == "vintage":
        add_rect(slide, Inches(4), Inches(4.5), Inches(5.3), Inches(0.04), palette["secondary"])
    elif not is_dark:
        add_rect(slide, Inches(4), Inches(4.5), Inches(5.3), Inches(0.04),
                 RGBColor(0xFF, 0xFF, 0xFF))

    subtitle = slide_data.get("content", [])
    if subtitle:
        sub_text = subtitle[0] if isinstance(subtitle, list) else subtitle
        sub_color = palette["subtext"] if (is_dark or style in ("swiss_modern", "vintage")) else RGBColor(0xCC, 0xDD, 0xEE)
        add_text_box(slide, sub_text,
                     Inches(1.5), Inches(4.8), Inches(10), Inches(0.8),
                     font_size=fonts["size_body"], color=sub_color,
                     align=PP_ALIGN.CENTER)
    return slide


def _apply_slide_bg(slide, palette):
    """Apply style-appropriate background to a content slide."""
    style = _current_style
    is_dark = style in ("bold_signal", "dark_botanical", "neon_cyber")
    is_notebooklm = style in ("minimalist", "editorial", "neo_retro")

    if is_dark:
        try:
            add_gradient_fill(slide, palette["gradient_start"], palette["gradient_end"])
        except Exception:
            set_bg(slide, palette["bg"])
    elif is_notebooklm:
        set_bg(slide, palette["bg"])
    else:
        set_bg(slide, palette["bg"])


def _add_title_bar(slide, title, palette):
    """Add style-appropriate title bar to content slides."""
    style = _current_style
    fonts = FONTS.get(style, FONTS["academic"])
    is_dark = style in ("bold_signal", "dark_botanical", "neon_cyber")
    is_notebooklm = style in ("minimalist", "editorial", "neo_retro")

    if style == "swiss_modern":
        # No full-width bar — just bold text + red accent line
        set_bg(slide, palette["bg"])
        add_text_box(slide, title,
                     Inches(0.6), Inches(0.2), Inches(12.1), Inches(1.0),
                     font_size=fonts["size_h2"], bold=True, color=palette["text"])
        add_rect(slide, Inches(0.6), Inches(1.15), Inches(2.5), Inches(0.05), palette["accent"])
    elif style == "notebook":
        set_bg(slide, palette["bg"])
        add_text_box(slide, title,
                     Inches(0.6), Inches(0.2), Inches(12.1), Inches(1.0),
                     font_size=fonts["size_h2"], bold=True, color=palette["text"])
        add_rect(slide, Inches(0.6), Inches(1.15), SLIDE_W - Inches(1.2), Inches(0.03),
                 palette["subtext"])
    elif style == "vintage":
        set_bg(slide, palette["bg"])
        add_text_box(slide, title,
                     Inches(0.8), Inches(0.2), Inches(11.7), Inches(1.0),
                     font_size=fonts["size_h2"], bold=True, color=palette["text"],
                     align=PP_ALIGN.CENTER)
        add_rect(slide, Inches(4), Inches(1.15), Inches(5.3), Inches(0.03), palette["secondary"])
    elif is_dark:
        # Dark themes: translucent title bar
        add_translucent_shape(slide, 1, Inches(0), Inches(0), SLIDE_W, Inches(1.2),
                              palette["primary"], opacity_pct=40)
        add_text_box(slide, title,
                     Inches(0.6), Inches(0.1), Inches(12.1), Inches(1.0),
                     font_size=fonts["size_h2"], bold=True, color=palette["text"])
    elif is_notebooklm:
        # NotebookLM styles: minimal title bar
        if style == "editorial":
            # Yellow accent line under title
            add_text_box(slide, title,
                         Inches(0.6), Inches(0.2), Inches(12.1), Inches(1.0),
                         font_size=fonts["size_h2"] + 4, bold=True, color=palette["text"])
            add_rect(slide, Inches(0.6), Inches(1.1), Inches(2), Inches(0.08), palette["primary"])
        elif style == "minimalist":
            # Super minimal - just text
            add_text_box(slide, title,
                         Inches(0.6), Inches(0.3), Inches(12.1), Inches(1.0),
                         font_size=fonts["size_h2"], bold=True, color=palette["text"])
            # Thin gray line
            add_rect(slide, Inches(0.6), Inches(1.15), Inches(12.1), Inches(0.02), palette["subtext"])
        elif style == "neo_retro":
            # Colorful accent
            add_text_box(slide, title,
                         Inches(0.6), Inches(0.2), Inches(12.1), Inches(1.0),
                         font_size=fonts["size_h2"], bold=True, color=palette["text"])
            add_rect(slide, Inches(0.6), Inches(1.1), Inches(3), Inches(0.06), palette["primary"])
    else:
        # Default: solid color title bar
        add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), palette["primary"])
        add_text_box(slide, title,
                     Inches(0.4), Inches(0.1), Inches(12.5), Inches(1.0),
                     font_size=fonts["size_h2"], bold=True,
                     color=RGBColor(0xFF, 0xFF, 0xFF))
        add_rect(slide, Inches(0), Inches(1.2), SLIDE_W, Inches(0.05), palette["secondary"])


def _enforce_density(content: list, max_bullets=6) -> list:
    """Enforce content density limits. Truncate to max_bullets."""
    if len(content) <= max_bullets:
        return content
    return content[:max_bullets]


def build_content_slide(prs, slide_data: dict, palette: dict):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _apply_slide_bg(slide, palette)

    title = slide_data.get("title", "")
    _add_title_bar(slide, title, palette)

    # Bullets with density control
    content = slide_data.get("content", [])
    content = _enforce_density(content, max_bullets=6)
    if content:
        txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        bullets = parse_bullets(content)
        add_bullets_to_tf(tf, bullets, palette, base_size=20)

    return slide


def build_two_col_slide(prs, slide_data: dict, palette: dict):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _apply_slide_bg(slide, palette)

    title = slide_data.get("title", "")
    _add_title_bar(slide, title, palette)

    # Divider
    add_rect(slide, Inches(6.55), Inches(1.5), Inches(0.04), Inches(5.6), palette["table_alt"])

    left = _enforce_density(slide_data.get("left", []), max_bullets=5)
    right = _enforce_density(slide_data.get("right", []), max_bullets=5)

    for col_content, left_pos in [(left, Inches(0.5)), (right, Inches(6.8))]:
        if col_content:
            txBox = slide.shapes.add_textbox(left_pos, Inches(1.5), Inches(5.9), Inches(5.6))
            tf = txBox.text_frame
            tf.word_wrap = True
            bullets = parse_bullets(col_content)
            add_bullets_to_tf(tf, bullets, palette, base_size=18)

    return slide


def build_table_slide(prs, slide_data: dict, palette: dict):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _apply_slide_bg(slide, palette)

    title = slide_data.get("title", "")
    _add_title_bar(slide, title, palette)

    table_data = slide_data.get("table", [])
    if not table_data:
        return slide

    rows = len(table_data)
    cols = max(len(r) for r in table_data)
    if rows == 0 or cols == 0:
        return slide

    # Calculate sizes
    tbl_width = Inches(12.3)
    tbl_height = min(Inches(5.5), Inches(0.5 * rows))
    left = Inches(0.5)
    top = Inches(1.4)

    table = slide.shapes.add_table(rows, cols, left, top, tbl_width, tbl_height).table

    # Apply enhanced table styling
    for r_idx, row in enumerate(table_data):
        for c_idx in range(cols):
            cell = table.cell(r_idx, c_idx)
            text = row[c_idx] if c_idx < len(row) else ""
            cell.text = str(text)
            tf = cell.text_frame

            # Enhanced header row styling
            if r_idx == 0:
                # Gradient-like effect with solid fill using accent
                cell.fill.solid()
                cell.fill.fore_color.rgb = palette["primary"]
                tf.paragraphs[0].font.size = Pt(16)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                # Center align header text
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            else:
                # Data row styling
                tf.paragraphs[0].font.size = Pt(14)
                tf.paragraphs[0].font.color.rgb = palette["text"]

            # Enhanced zebra striping with background_tint color
            if r_idx > 0 and r_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = palette["background_tint"]

            # Enhanced border styling
            border = cell.border
            border.top.color.rgb = palette["subtext"]
            border.bottom.color.rgb = palette["subtext"]
            border.left.color.rgb = palette["subtext"]
            border.right.color.rgb = palette["subtext"]
            border.top.width = Pt(1)
            border.bottom.width = Pt(1)
            border.left.width = Pt(1)
            border.right.width = Pt(1)

    return slide


def build_three_col_slide(prs, slide_data: dict, palette: dict):
    """Three equal columns for comparisons."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _apply_slide_bg(slide, palette)

    title = slide_data.get("title", "")
    _add_title_bar(slide, title, palette)

    # Three equal columns
    col_width = Inches(4.0)
    positions = [Inches(0.4), Inches(4.5), Inches(8.6)]
    dividers = [Inches(4.4), Inches(8.5)]

    # Vertical dividers
    for div_pos in dividers:
        add_rect(slide, div_pos, Inches(1.4), Inches(0.05), Inches(5.6), palette["table_alt"])

    for i in range(3):
        col_key = f"col{i+1}"
        col_content = slide_data.get(col_key, slide_data.get("content", []))
        if col_content:
            txBox = slide.shapes.add_textbox(positions[i], Inches(1.4), col_width, Inches(5.6))
            tf = txBox.text_frame
            tf.word_wrap = True
            bullets = parse_bullets(col_content if isinstance(col_content, list) else [col_content])
            add_bullets_to_tf(tf, bullets, palette, base_size=16)

    return slide


def build_left_sidebar_slide(prs, slide_data: dict, palette: dict):
    """Content area with left sidebar for notes/definitions."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _apply_slide_bg(slide, palette)

    title = slide_data.get("title", "")
    _add_title_bar(slide, title, palette)

    # Sidebar background (left 30%)
    add_rect(slide, Inches(0), Inches(1.25), Inches(3.9), Inches(6.0), palette["table_alt"])

    # Divider
    add_rect(slide, Inches(3.9), Inches(1.4), Inches(0.05), Inches(5.6), palette["secondary"])

    # Sidebar content (notes/definitions)
    sidebar = slide_data.get("sidebar", slide_data.get("left", []))
    if sidebar:
        txBox = slide.shapes.add_textbox(Inches(0.3), Inches(1.5), Inches(3.5), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        bullets = parse_bullets(sidebar if isinstance(sidebar, list) else [sidebar])
        add_bullets_to_tf(tf, bullets, palette, base_size=14)

    # Main content
    content = slide_data.get("content", [])
    if content:
        txBox = slide.shapes.add_textbox(Inches(4.1), Inches(1.5), Inches(8.7), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        bullets = parse_bullets(content)
        add_bullets_to_tf(tf, bullets, palette, base_size=18)

    return slide


def build_right_sidebar_slide(prs, slide_data: dict, palette: dict):
    """Content area with right sidebar for stats/highlights."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _apply_slide_bg(slide, palette)

    title = slide_data.get("title", "")
    _add_title_bar(slide, title, palette)

    # Sidebar background (right 30%)
    add_rect(slide, Inches(9.4), Inches(1.25), Inches(3.9), Inches(6.0), palette["table_alt"])

    # Divider
    add_rect(slide, Inches(9.35), Inches(1.4), Inches(0.05), Inches(5.6), palette["secondary"])

    # Main content
    content = slide_data.get("content", [])
    if content:
        txBox = slide.shapes.add_textbox(Inches(0.4), Inches(1.5), Inches(8.7), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        bullets = parse_bullets(content)
        add_bullets_to_tf(tf, bullets, palette, base_size=18)

    # Sidebar content (stats/highlights)
    sidebar = slide_data.get("sidebar", slide_data.get("right", []))
    if sidebar:
        txBox = slide.shapes.add_textbox(Inches(9.5), Inches(1.5), Inches(3.7), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        bullets = parse_bullets(sidebar if isinstance(sidebar, list) else [sidebar])
        add_bullets_to_tf(tf, bullets, palette, base_size=14)

    return slide


def build_quote_slide(prs, slide_data: dict, palette: dict):
    """Large quote with attribution."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _apply_slide_bg(slide, palette)

    # Decorative accent bar on left
    add_rect(slide, Inches(0), Inches(0), Inches(0.6), SLIDE_H, palette["accent"])

    # Quote content
    quote_text = slide_data.get("content", [""])[0] if slide_data.get("content") else ""
    add_text_box(slide, f'"{quote_text}"',
                 Inches(1.2), Inches(2.0), Inches(11.5), Inches(3.5),
                 font_size=32, color=palette["primary"], align=PP_ALIGN.CENTER)

    # Attribution
    attribution = slide_data.get("attribution", slide_data.get("title", ""))
    if attribution:
        add_text_box(slide, f"— {attribution}",
                     Inches(1.2), Inches(5.5), Inches(11.5), Inches(0.8),
                     font_size=20, bold=True, color=palette["subtext"], align=PP_ALIGN.CENTER)

    # Accent line below
    add_rect(slide, Inches(4), Inches(6.2), Inches(5.3), Inches(0.08), palette["accent"])

    return slide


def build_center_focus_slide(prs, slide_data: dict, palette: dict):
    """Centered content with accent border."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _apply_slide_bg(slide, palette)

    # Outer accent border
    border_thickness = Inches(0.15)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, border_thickness, palette["primary"])  # Top
    add_rect(slide, Inches(0), SLIDE_H - border_thickness, SLIDE_W, border_thickness, palette["primary"])  # Bottom
    add_rect(slide, Inches(0), Inches(0), border_thickness, SLIDE_H, palette["primary"])  # Left
    add_rect(slide, SLIDE_W - border_thickness, Inches(0), border_thickness, SLIDE_H, palette["primary"])  # Right

    # Title
    title = slide_data.get("title", "")
    add_text_box(slide, title,
                 Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.2),
                 font_size=32, bold=True, color=palette["primary"], align=PP_ALIGN.CENTER)

    # Centered content
    content = slide_data.get("content", [])
    if content:
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(10.3), Inches(4.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        bullets = parse_bullets(content)
        add_bullets_to_tf(tf, bullets, palette, base_size=20)
        # Center align all paragraphs
        for p in tf.paragraphs:
            p.alignment = PP_ALIGN.CENTER

    return slide


def build_comparison_slide(prs, slide_data: dict, palette: dict):
    """Before/after or Option A/B comparison matrix."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _apply_slide_bg(slide, palette)

    title = slide_data.get("title", "")
    _add_title_bar(slide, title, palette)

    # Get column labels (default: "Option A", "Option B")
    left_label = slide_data.get("left_label", "Option A")
    right_label = slide_data.get("right_label", "Option B")

    # Column headers with different backgrounds
    add_rect(slide, Inches(0.4), Inches(1.4), Inches(6.2), Inches(0.6), palette["secondary"])
    add_rect(slide, Inches(6.7), Inches(1.4), Inches(6.2), Inches(0.6), palette["primary"])

    add_text_box(slide, left_label,
                 Inches(0.4), Inches(1.45), Inches(6.2), Inches(0.5),
                 font_size=20, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    add_text_box(slide, right_label,
                 Inches(6.7), Inches(1.45), Inches(6.2), Inches(0.5),
                 font_size=20, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

    # Content columns
    left = slide_data.get("left", slide_data.get("content", []))
    right = slide_data.get("right", [])

    for col_content, left_pos in [(left, Inches(0.4)), (right, Inches(6.7))]:
        if col_content:
            txBox = slide.shapes.add_textbox(left_pos, Inches(2.1), Inches(6.2), Inches(5.0))
            tf = txBox.text_frame
            tf.word_wrap = True
            bullets = parse_bullets(col_content if isinstance(col_content, list) else [col_content])
            add_bullets_to_tf(tf, bullets, palette, base_size=16)

    return slide


def build_process_slide(prs, slide_data: dict, palette: dict):
    """Step-by-step process flow with numbered steps."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _apply_slide_bg(slide, palette)

    title = slide_data.get("title", "")
    _add_title_bar(slide, title, palette)

    # Process steps - horizontal layout
    steps = slide_data.get("steps", slide_data.get("content", []))
    if not steps:
        return slide

    num_steps = min(len(steps), 5)  # Max 5 steps horizontally
    step_width = Inches(12.0 / num_steps)

    for i, step in enumerate(steps[:5]):
        x_pos = Inches(0.5 + i * (12.0 / num_steps))

        # Step circle with number
        circle = slide.shapes.add_shape(
            9,  # MSO_SHAPE_TYPE.OVAL
            x_pos, Inches(1.5), Inches(0.7), Inches(0.7)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = palette["accent"]
        circle.line.color.rgb = palette["primary"]
        circle.line.width = Pt(2)

        # Step number
        txBox = slide.shapes.add_textbox(x_pos, Inches(1.6), Inches(0.7), Inches(0.5))
        tf = txBox.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = str(i + 1)
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # Step text below
        txBox = slide.shapes.add_textbox(x_pos, Inches(2.4), step_width, Inches(3.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = step if isinstance(step, str) else step[0] if step else ""
        run.font.size = Pt(14)
        run.font.color.rgb = palette["text"]

        # Arrow to next step (except last)
        if i < num_steps - 1:
            add_rect(slide, x_pos + Inches(0.7), Inches(1.8), Inches(0.3), Inches(0.05), palette["subtext"])

    return slide


def build_chart_slide(prs, slide_data: dict, palette: dict):
    """Slide with generated chart image."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _apply_slide_bg(slide, palette)

    title = slide_data.get("title", "")
    _add_title_bar(slide, title, palette)

    # Get chart specification
    chart_spec = slide_data.get("chart", {})
    if not chart_spec:
        return slide

    # Generate chart image
    chart_buf = create_chart_from_spec(chart_spec, palette)
    if chart_buf:
        # Save chart to temp file and insert
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(chart_buf.getvalue())
            tmp_path = tmp.name

        try:
            # Add chart image centered
            img_width = Inches(10)
            img_height = Inches(5)
            left = (SLIDE_W - img_width) / 2
            top = Inches(1.5)
            slide.shapes.add_picture(tmp_path, left, top, width=img_width, height=img_height)
        finally:
            os.unlink(tmp_path)

    # Add caption/notes if provided
    caption = slide_data.get("caption", slide_data.get("content", [""])[0] if slide_data.get("content") else "")
    if caption:
        add_text_box(slide, caption,
                     Inches(0.5), Inches(6.6), Inches(12.3), Inches(1.0),
                     font_size=14, color=palette["subtext"])

    return slide


BUILDERS = {
    "title":          build_title_slide,
    "section":        build_section_slide,
    "content":        build_content_slide,
    "two_col":        build_two_col_slide,
    "table":          build_table_slide,
    "three_col":      build_three_col_slide,
    "left_sidebar":   build_left_sidebar_slide,
    "right_sidebar":  build_right_sidebar_slide,
    "quote":          build_quote_slide,
    "center_focus":   build_center_focus_slide,
    "comparison":     build_comparison_slide,
    "process":        build_process_slide,
    "image_left":     build_image_left_slide,
    "image_right":    build_image_right_slide,
    "image_top":      build_image_top_slide,
    "chart":          build_chart_slide,
    # NotebookLM-inspired layouts
    "big_number":     build_big_number_slide,
    "split_panel":    build_split_panel_slide,
    "timeline":       build_timeline_slide,
    "card_grid":      build_card_grid_slide,
}


def add_speaker_notes(slide, notes_text: str):
    if not notes_text:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


def generate_style_previews(output_dir="/workspace/group/previews"):
    """Generate preview images for each style - DEPRECATED: use generate_all_styles_preview instead."""
    try:
        os.makedirs(output_dir, exist_ok=True)

        for style_name in PALETTES.keys():
            # Create a sample presentation
            prs = Presentation()
            prs.slide_width = SLIDE_W
            prs.slide_height = SLIDE_H

            palette = get_palette(style_name)

            # Create a sample slide with title and content
            slide_layout = prs.slide_layouts[6]  # blank
            slide = prs.slides.add_slide(slide_layout)
            set_bg(slide, palette["bg"])

            # Add decorative elements
            add_accent_bar(slide, position="left", palette=palette)

            # Title
            add_text_box(slide, f"{style_name.title()} Style",
                         Inches(1.5), Inches(2.0), Inches(10), Inches(0.8),
                         font_size=36, bold=True, color=palette["primary"])

            # Sample content
            add_text_box(slide, f"This is a sample {style_name} presentation style.",
                         Inches(1.5), Inches(3.0), Inches(10), Inches(0.6),
                         font_size=18, color=palette["text"])

            # Add decorative divider
            add_rect(slide, Inches(3), Inches(3.8), Inches(7), Inches(0.08), palette["secondary"])

            # Sample bullets
            txBox = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10), Inches(2.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            bullets = [
                (0, "Feature 1: Professional typography"),
                (0, "Feature 2: Coordinated color palette"),
                (0, "Feature 3: Visual depth and accents"),
            ]
            add_bullets_to_tf(tf, bullets, palette, base_size=16)

            # Save preview
            preview_path = os.path.join(output_dir, f"{style_name}_preview.png")
            prs.save(preview_path)
            print(f"Generated preview: {preview_path}")

    except Exception as e:
        print(f"WARNING: Could not generate previews: {e}", file=sys.stderr)


def generate_all_styles_preview(output_path="/workspace/group/all_styles_preview.pptx"):
    """Generate a single PPTX with sample slides for ALL 13 styles.
    This allows users to visually compare and choose their preferred style.
    Each style gets 1 slide showing:
    - Style name and description
    - Sample content with bullets
    - Style-specific decorations and colors
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    global _current_style
    original_style = _current_style

    # Style descriptions
    style_descriptions = {
        "academic": "经典学术 · 深蓝+钢蓝 · 适合开题答辩、学术汇报",
        "clinical": "临床教学 · 蓝色+红色强调 · 适合科室培训、病例讨论",
        "patient": "患者教育 · 温暖青色 · 适合健康宣教、科普",
        "concise": "简洁专业 · 灰蓝极简 · 适合快速分享、业务汇报",
        "bold_signal": "科技感 · 暗底+电蓝+粉红 · 适合创新展示",
        "dark_botanical": "自然优雅 · 森林绿+鼠尾草 · 适合高端演示",
        "notebook": "温馨笔记 · 米色+陶土橙 · 适合培训记录",
        "neon_cyber": "赛博朋克 · 深紫+霓虹绿 · 适合技术展示",
        "swiss_modern": "瑞士设计 · 白底+红+黑 · 适合简洁专业",
        "vintage": "复古优雅 · 奶油色+深绿+金色 · 适合文化主题",
        "minimalist": "极简主义 · 浅灰+黑字+黄点缀 · 适合高端展示",
        "editorial": "杂志风格 · 白底+黄色主色块+红强调 · 适合视觉冲击",
        "neo_retro": "复古科技 · 奶油格子纸+粉/青/黄 · 适合创意展示",
    }

    for style_name in PALETTES.keys():
        _current_style = style_name
        palette = get_palette(style_name)
        fonts = FONTS.get(style_name, FONTS["academic"])
        is_dark = style_name in ("bold_signal", "dark_botanical", "neon_cyber")

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Apply background
        if is_dark:
            try:
                add_gradient_fill(slide, palette["gradient_start"], palette["gradient_end"])
            except Exception:
                set_bg(slide, palette["bg"])
        else:
            set_bg(slide, palette["bg"])

        # Style name as title
        style_display = style_name.replace("_", " ").title()
        add_text_box(slide, style_display,
                     Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.2),
                     font_size=fonts["size_h1"], bold=True, color=palette.get("text", palette["primary"]))

        # Description
        desc = style_descriptions.get(style_name, "")
        add_text_box(slide, desc,
                     Inches(0.8), Inches(2.8), Inches(11.5), Inches(0.8),
                     font_size=fonts["size_body"], color=palette["subtext"])

        # Sample content
        sample_content = [
            "清晰的视觉层次",
            "专注核心信息",
            "独特的品牌识别",
        ]

        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(6), Inches(2.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        bullets = parse_bullets(sample_content)
        add_bullets_to_tf(tf, bullets, palette, base_size=18)

        # Apply style-specific decorations
        add_decorations(slide, "content", palette)

        # Add style number in corner
        style_num = list(PALETTES.keys()).index(style_name) + 1
        add_text_box(slide, f"{style_num:02d}" if style_num < 10 else f"{style_num}",
                     SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.6), Inches(1), Inches(0.4),
                     font_size=12, color=palette["subtext"], align=PP_ALIGN.RIGHT)

    # Restore original style
    _current_style = original_style

    # Ensure output directory exists
    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    prs.save(output_path)
    return output_path


def auto_generate_images(slides: list, style: str, output_dir: str) -> list:
    """Auto-generate images for slides that could benefit from visuals.
    Modifies slides in-place, adding image paths and switching layouts."""
    try:
        sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
        from gemini_image_gen import generate_image
    except ImportError:
        print("WARNING: gemini_image_gen not available, skipping auto-images", file=sys.stderr)
        return slides

    img_dir = os.path.join(output_dir, "generated_images")
    os.makedirs(img_dir, exist_ok=True)

    image_count = 0
    max_images = 4  # Don't overdo it

    for i, slide_data in enumerate(slides):
        if image_count >= max_images:
            break

        layout = slide_data.get("layout", "content")
        title = slide_data.get("title", "")

        # Skip layouts that already have images or aren't suitable
        if layout in ("title", "table", "chart", "image_left", "image_right", "image_top"):
            continue
        if slide_data.get("image"):
            continue

        # Generate image for section slides (chapter covers)
        if layout == "section" and title:
            img_path = os.path.join(img_dir, f"section_{i}.png")
            try:
                generate_image({
                    "prompt": f"Abstract, minimalist medical illustration representing: {title}. Clean lines, professional, suitable for presentation background.",
                    "output_path": img_path,
                    "style": "medical",
                })
                # Don't change section layout, but note the image for potential use
                image_count += 1
            except Exception as e:
                print(f"WARNING: Image generation failed for slide {i}: {e}", file=sys.stderr)
            continue

        # For content slides, generate image and switch to image_right layout
        if layout == "content" and title and image_count < max_images:
            # Only generate for every other content slide to avoid monotony
            content_slides_before = sum(1 for s in slides[:i] if s.get("layout") == "content")
            if content_slides_before % 3 != 0:
                continue

            img_path = os.path.join(img_dir, f"content_{i}.png")
            try:
                generate_image({
                    "prompt": f"Professional medical illustration for a presentation slide about: {title}. Clean, modern style, white background, suitable for medical presentation.",
                    "output_path": img_path,
                    "style": "medical",
                })
                if os.path.exists(img_path):
                    slide_data["layout"] = "image_right"
                    slide_data["image"] = {"path": img_path, "style": "plain"}
                    image_count += 1
            except Exception as e:
                print(f"WARNING: Image generation failed for slide {i}: {e}", file=sys.stderr)

    print(f"Auto-generated {image_count} images", file=sys.stderr)
    return slides


def auto_generate_images_aggressive(slides: list, style: str, output_dir: str) -> list:
    """Aggressively generate images for ALL content slides.
    This is the 'full throttle' version that ensures every slide has visuals."""
    try:
        sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
        from gemini_image_gen import generate_image
    except ImportError:
        print("WARNING: gemini_image_gen not available, skipping auto-images", file=sys.stderr)
        return slides

    img_dir = os.path.join(output_dir, "generated_images")
    os.makedirs(img_dir, exist_ok=True)

    image_count = 0
    max_images = 20  # Generate many more images

    # Style-specific image prompt modifiers
    style_prompts = {
        "academic": "Professional medical journal illustration style, clean and scientific",
        "clinical": "Clinical medical diagram, clean white background, professional healthcare style",
        "patient": "Friendly patient education illustration, warm colors, approachable style",
        "concise": "Minimalist medical icon or diagram, simple and clear",
        "bold_signal": "Bold graphic medical illustration, high contrast, modern design",
        "dark_botanical": "Artistic medical illustration with natural elements, sophisticated",
        "notebook": "Hand-drawn style medical sketch, notebook paper aesthetic",
        "neon_cyber": "Futuristic medical visualization, glowing accents",
        "swiss_modern": "Clean Swiss design medical diagram, grid-based, minimal",
        "vintage": "Vintage medical illustration style, classic textbook aesthetic",
        "minimalist": "Ultra-minimalist medical icon, maximum whitespace",
        "editorial": "Magazine-style medical illustration, bold and editorial",
        "neo_retro": "Modern retro medical graphic, colorful and playful",
    }
    style_modifier = style_prompts.get(style, style_prompts["clinical"])

    for i, slide_data in enumerate(slides):
        if image_count >= max_images:
            break

        layout = slide_data.get("layout", "content")
        title = slide_data.get("title", "")
        content = slide_data.get("content", [])

        # Skip layouts that already have images or are tables
        if layout in ("title", "table", "chart", "image_left", "image_right", "image_top"):
            continue
        if slide_data.get("image"):
            continue

        # Skip section slides - they work fine without images
        if layout == "section":
            continue

        # Generate image for content slides - ALL of them
        if layout in ("content", "two_col", "three_col", "comparison", "process",
                      "big_number", "split_panel", "timeline", "card_grid"):
            # Build context-aware prompt from slide content
            content_summary = title
            if content:
                # Take first 2 bullet points for context
                content_summary = f"{title}: {', '.join(content[:2])}"

            img_path = os.path.join(img_dir, f"slide_{i:02d}.png")

            # Enhanced prompt with style and content context
            prompt = (
                f"{style_modifier}. "
                f"Illustration for medical presentation slide about: {content_summary}. "
                f"Clean, professional, suitable for PowerPoint. "
                f"No text, no labels, pure visual illustration. "
                f"Aspect ratio suitable for slide sidebar placement."
            )

            try:
                generate_image({
                    "prompt": prompt,
                    "output_path": img_path,
                    "style": "medical",
                })
                if os.path.exists(img_path):
                    # Switch to image_right layout to display the generated image
                    slide_data["layout"] = "image_right"
                    slide_data["image"] = {"path": img_path, "style": "plain"}
                    image_count += 1
                    print(f"Generated image for slide {i}: {title[:30]}...", file=sys.stderr)
            except Exception as e:
                print(f"WARNING: Image generation failed for slide {i}: {e}", file=sys.stderr)
                # Continue without image - don't fail the whole presentation

    print(f"Aggressive auto-generation complete: {image_count} images generated", file=sys.stderr)
    return slides


def generate_style_preview(output_path: str = "/workspace/group/style_preview.pptx") -> str:
    """Generate a single PPTX with 3 sample slides showing different styles.
    This allows users to visually compare and choose their preferred style."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Select 3 representative styles for preview
    preview_styles = ["minimalist", "editorial", "neo_retro"]

    global _current_style
    original_style = _current_style

    for i, style_name in enumerate(preview_styles):
        _current_style = style_name
        palette = get_palette(style_name)
        fonts = FONTS.get(style_name, FONTS["academic"])

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        is_dark = style_name in ("bold_signal", "dark_botanical", "neon_cyber")

        # Apply background
        if is_dark:
            try:
                add_gradient_fill(slide, palette["gradient_start"], palette["gradient_end"])
            except Exception:
                set_bg(slide, palette["bg"])
        else:
            set_bg(slide, palette["bg"])

        # Style name as title
        style_display = {
            "minimalist": "Minimalist",
            "editorial": "Editorial",
            "neo_retro": "Neo Retro",
            "academic": "Academic",
            "clinical": "Clinical",
            "patient": "Patient",
            "concise": "Concise",
            "bold_signal": "Bold Signal",
            "dark_botanical": "Dark Botanical",
            "notebook": "Notebook",
            "neon_cyber": "Neon Cyber",
            "swiss_modern": "Swiss Modern",
            "vintage": "Vintage",
        }.get(style_name, style_name.title())

        add_text_box(slide, style_display,
                     Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.5),
                     font_size=fonts["size_h1"], bold=True, color=palette["text"])

        # Sample subtitle
        subtitle = {
            "minimalist": "极简 · 留白 · 聚焦",
            "editorial": "杂志 · 黄黑 · 冲击",
            "neo_retro": "复古 · 多彩 · 科技",
        }.get(style_name, "专业医学风格")

        add_text_box(slide, subtitle,
                     Inches(0.8), Inches(3.5), Inches(11.5), Inches(0.6),
                     font_size=fonts["size_body"], color=palette["subtext"])

        # Sample content bullets
        sample_content = [
            "清晰的视觉层次",
            "专注核心信息",
            "独特的品牌识别",
        ]

        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(6), Inches(2.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        bullets = parse_bullets(sample_content)
        add_bullets_to_tf(tf, bullets, palette, base_size=18)

        # Apply style-specific decorations
        add_decorations(slide, "content", palette)

        # Add style indicator in corner
        add_text_box(slide, f"0{i+1}/3",
                     SLIDE_W - Inches(1.5), SLIDE_H - Inches(0.6), Inches(1), Inches(0.4),
                     font_size=12, color=palette["subtext"], align=PP_ALIGN.RIGHT)

    # Restore original style
    _current_style = original_style

    prs.save(output_path)
    return output_path


def generate(data: dict) -> str:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    global _current_style
    # Default to 'clinical' for medical presentations - NEVER use neon_cyber as default!
    style = data.get("style", "clinical")
    _current_style = style
    palette = get_palette(style)
    output_path = data.get("output_path", "/workspace/group/presentation.pptx")

    # Ensure output directory exists
    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    # ALWAYS auto-generate images if ZENMUX_API_KEY is available
    # Check env and secrets file
    slides = data.get("slides", [])
    zenmux_key = os.environ.get("ZENMUX_API_KEY", "")
    if not zenmux_key:
        # Try reading from secrets file written by agent runner
        try:
            with open("/tmp/.zenmux_key", "r") as f:
                zenmux_key = f.read().strip()
                if zenmux_key:
                    os.environ["ZENMUX_API_KEY"] = zenmux_key
        except FileNotFoundError:
            pass

    # FORCE image generation for ALL content slides (not just 1/3)
    if zenmux_key and not data.get("skip_auto_images"):
        try:
            slides = auto_generate_images_aggressive(slides, style, os.path.dirname(output_path))
        except Exception as e:
            print(f"WARNING: Auto image generation failed: {e}", file=sys.stderr)

    # Always add title slide first
    title_slide = build_title_slide(prs, data, palette)

    for slide_data in slides:
        layout = slide_data.get("layout", "content")
        if layout == "title":
            continue
        builder = BUILDERS.get(layout, build_content_slide)
        slide = builder(prs, slide_data, palette)

        # Add decorations
        add_decorations(slide, layout, palette)

        notes = slide_data.get("notes", "")
        if notes:
            add_speaker_notes(slide, notes)

    prs.save(output_path)
    return output_path


def main():
    if len(sys.argv) > 1:
        arg = sys.argv[1]

        # Check for preview mode
        if arg == "--preview" or arg == "preview":
            # Optional: specify output path
            output_path = "/workspace/group/all_styles_preview.pptx"

            # Parse additional arguments
            args = sys.argv[2:]
            i = 0
            while i < len(args):
                if args[i] == "--output" and i + 1 < len(args):
                    output_path = args[i + 1]
                    i += 2
                else:
                    i += 1

            path = generate_all_styles_preview(output_path)
            print(path)
            return

        # Normal JSON input mode
        raw = arg
    else:
        raw = sys.stdin.read()

    try:
        data = json.loads(raw)
    except json.JSONDecodeError as e:
        print(f"ERROR: Invalid JSON input: {e}", file=sys.stderr)
        sys.exit(1)

    path = generate(data)
    print(path)


if __name__ == "__main__":
    main()
